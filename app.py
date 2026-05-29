import streamlit as st
import PyPDF2
import docx
import requests
from bs4 import BeautifulSoup
from io import StringIO
import base64
from collections import Counter
import re

# Page config
st.set_page_config(page_title="🤖 Ayushi's AI Analyzer", layout="wide", page_icon="🤖")

# Custom CSS
st.markdown("""
<style>
    .main-header {font-size: 3rem; font-weight: bold; color: #FF4B4B;}
    .sub-header {font-size: 1.5rem; color: #FF6347;}
    .stTextInput>div>div>input {color: #4CAF50;}
</style>
""", unsafe_allow_html=True)

# Title
st.markdown('<p class="main-header">🤖 Ayushi\'s AI Analyzer</p>', unsafe_allow_html=True)
st.markdown("Upload documents OR paste website links — I'll analyze them like ChatGPT!")

# Sidebar
st.sidebar.title("📥 Input Method")
input_method = st.sidebar.radio("Choose input type:", ["📁 Upload File", "🌐 Website URL"])

# Initialize session state
if 'content_text' not in st.session_state:
    st.session_state.content_text = ""
if 'source_name' not in st.session_state:
    st.session_state.source_name = ""
if 'content_type' not in st.session_state:
    st.session_state.content_type = ""

# File processing functions
def extract_pdf_text(file):
    try:
        pdf_reader = PyPDF2.PdfReader(file)
        text = ""
        for page_num, page in enumerate(pdf_reader.pages, 1):
            text += f"\n--- Page {page_num} ---\n"
            text += page.extract_text() or ""
        return text
    except Exception as e:
        return f"Error: {e}"

def extract_docx_text(file):
    doc = docx.Document(file)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

def extract_txt_text(file):
    return file.getvalue().decode("utf-8")

def extract_pptx_text(file):
    try:
        from pptx import Presentation
        prs = Presentation(file)
        text = ""
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"\n--- Slide {slide_num} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except ImportError:
        return "Please install python-pptx: pip install python-pptx"
    except Exception as e:
        return f"Error: {e}"

def scrape_website(url):
    try:
        headers = {'User-Agent': 'Mozilla/5.0'}
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        
        soup = BeautifulSoup(response.content, 'html.parser')
        
        # Remove script and style elements
        for script in soup(['script', 'style', 'nav', 'footer', 'header']):
            script.decompose()
        
        text = soup.get_text(separator='\n', strip=True)
        return text
    except Exception as e:
        return f"Error scraping website: {e}"

# FILE UPLOAD SECTION
if input_method == "📁 Upload File":
    uploaded_file = st.sidebar.file_uploader(
        "Choose a file", 
        type=['txt', 'pdf', 'docx', 'pptx', 'md', 'csv', 'json']
    )
    
    if uploaded_file is not None:
        file_type = uploaded_file.name.split('.')[-1].lower()
        st.sidebar.success(f"✅ Uploaded: {uploaded_file.name}")
        
        with st.spinner("📖 Reading your file..."):
            if file_type == 'pdf':
                st.session_state.content_text = extract_pdf_text(uploaded_file)
            elif file_type == 'docx':
                st.session_state.content_text = extract_docx_text(uploaded_file)
            elif file_type == 'txt':
                st.session_state.content_text = extract_txt_text(uploaded_file)
            elif file_type == 'pptx':
                st.session_state.content_text = extract_pptx_text(uploaded_file)
            elif file_type in ['md', 'csv', 'json']:
                st.session_state.content_text = extract_txt_text(uploaded_file)
            
            st.session_state.source_name = uploaded_file.name
            st.session_state.content_type = f"File ({file_type.upper()})"

# WEBSITE URL SECTION
elif input_method == "🌐 Website URL":
    url = st.sidebar.text_input("Enter website URL:", placeholder="https://example.com")
    
    if url:
        if not url.startswith(('http://', 'https://')):
            url = 'https://' + url
        
        with st.spinner("🌐 Scraping website..."):
            st.session_state.content_text = scrape_website(url)
            st.session_state.source_name = url
            st.session_state.content_type = "Website"

# MAIN CONTENT AREA
if st.session_state.content_text:
    # Show source info
    col1, col2, col3 = st.columns([3, 1, 1])
    with col1:
        st.markdown(f"### 📄 {st.session_state.source_name}")
        st.caption(f"Type: {st.session_state.content_type}")
    with col2:
        word_count = len(st.session_state.content_text.split())
        st.metric("Words", f"{word_count:,}")
    with col3:
        char_count = len(st.session_state.content_text)
        st.metric("Characters", f"{char_count:,}")
    
    st.divider()
    
    # Main interaction area
    st.markdown("### 🎯 What would you like me to do?")
    
    # Quick action buttons
    col1, col2, col3, col4, col5 = st.columns(5)
    with col1:
        if st.button("📝 Summarize", use_container_width=True):
            st.session_state.current_task = "summarize"
    with col2:
        if st.button("🔑 Key Points", use_container_width=True):
            st.session_state.current_task = "keypoints"
    with col3:
        if st.button("📊 Analyze", use_container_width=True):
            st.session_state.current_task = "analyze"
    with col4:
        if st.button("🔍 Search", use_container_width=True):
            st.session_state.current_task = "search"
    with col5:
        if st.button("📖 Read All", use_container_width=True):
            st.session_state.current_task = "read"
    
    if 'current_task' not in st.session_state:
        st.session_state.current_task = "summarize"
    
    st.divider()
    
    # TASK EXECUTION
    if st.session_state.current_task == "summarize":
        st.markdown("### 📝 Summary")
        
        # Auto-summary
        paragraphs = [p.strip() for p in st.session_state.content_text.split('\n\n') if len(p.strip()) > 100]
        
        if paragraphs:
            st.markdown("**Executive Summary:**")
            st.info(f"📌 This document contains {len(paragraphs)} main sections.")
            
            st.markdown("**Introduction:**")
            st.write(paragraphs[0][:500] + "..." if len(paragraphs[0]) > 500 else paragraphs[0])
            
            if len(paragraphs) > 1:
                st.markdown("**Conclusion:**")
                st.write(paragraphs[-1][:500] + "..." if len(paragraphs[-1]) > 500 else paragraphs[-1])
        
        # Extractive summary (important sentences)
        st.markdown("---")
        st.markdown("**Key Sentences:**")
        sentences = re.split(r'[.!?]+', st.session_state.content_text)
        sentences = [s.strip() for s in sentences if len(s.strip()) > 50]
        
        # Score sentences by word frequency
        words = st.session_state.content_text.lower().split()
        word_freq = Counter(words)
        
        scored_sentences = []
        for sent in sentences[:20]:
            score = sum(word_freq.get(word.lower(), 0) for word in sent.split())
            scored_sentences.append((sent, score))
        
        scored_sentences.sort(key=lambda x: x[1], reverse=True)
        
        for i, (sent, _) in enumerate(scored_sentences[:5], 1):
            st.write(f"{i}. {sent}.")
    
    elif st.session_state.current_task == "keypoints":
        st.markdown("### 🔑 Key Points Extracted")
        
        # Extract bullet points, numbered lists, important lines
        lines = st.session_state.content_text.split('\n')
        key_points = []
        
        for line in lines:
            line = line.strip()
            # Look for bullet points, numbered items, or important statements
            if any([
                line.startswith(('•', '-', '*', '▪', '▫', '1.', '2.', '3.', 'a.', 'b.', 'c.')),
                re.match(r'^[A-Z][^.!?]{50,150}[.!?]$', line),  # Capitalized sentences
                len(line.split()) < 15 and any(word in line.lower() for word in ['important', 'key', 'main', 'critical'])
            ]):
                if line and len(line) > 20:
                    key_points.append(line)
        
        if key_points:
            for point in key_points[:15]:
                st.markdown(f"• {point}")
        else:
            # Fallback: extract important sentences
            st.markdown("**Important Statements:**")
            sentences = re.split(r'[.!?]+', st.session_state.content_text)
            important = [s.strip() for s in sentences if 50 < len(s.strip()) < 200]
            for sent in important[:10]:
                st.markdown(f"• {sent}.")
    
    elif st.session_state.current_task == "analyze":
        st.markdown("### 📊 Content Analysis")
        
        col1, col2 = st.columns(2)
        
        with col1:
            st.markdown("**📈 Statistics**")
            text = st.session_state.content_text
            
            st.metric("Total Words", f"{len(text.split()):,}")
            st.metric("Total Characters", f"{len(text):,}")
            st.metric("Total Sentences", f"{len(re.split(r'[.!?]+', text)):,}")
            st.metric("Total Paragraphs", f"{len([p for p in text.split('\n\n') if p.strip()]):,}")
            st.metric("Avg Words per Sentence", f"{len(text.split()) / max(len(re.split(r'[.!?]+', text)), 1):.1f}")
        
        with col2:
            st.markdown("**🔤 Top 10 Keywords**")
            words = text.lower().split()
            # Remove stop words and short words
            stop_words = {'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 'is', 'are', 'was', 'were', 'be', 'been', 'being', 'have', 'has', 'had', 'do', 'does', 'did', 'will', 'would', 'could', 'should', 'may', 'might', 'must', 'shall', 'can', 'need', 'dare', 'ought', 'used', 'it', 'its', 'this', 'that', 'these', 'those', 'i', 'you', 'he', 'she', 'we', 'they', 'what', 'which', 'who', 'when', 'where', 'why', 'how', 'all', 'each', 'every', 'both', 'few', 'more', 'most', 'other', 'some', 'such', 'no', 'nor', 'not', 'only', 'own', 'same', 'so', 'than', 'too', 'very', 'just', 'as'}
            
            filtered_words = [w for w in words if w not in stop_words and len(w) > 3 and w.isalpha()]
            word_freq = Counter(filtered_words).most_common(10)
            
            for word, count in word_freq:
                st.write(f"• **{word}**: {count} times")
        
        st.divider()
        
        st.markdown("**📋 Content Type Detection**")
        text_lower = text.lower()
        
        content_types = []
        if any(word in text_lower for word in ['figure', 'table', 'chart', 'graph']):
            content_types.append("📊 Contains data/visualizations")
        if any(word in text_lower for word in ['introduction', 'conclusion', 'abstract']):
            content_types.append("📝 Academic/research paper")
        if any(word in text_lower for word in ['code', 'function', 'variable', 'program']):
            content_types.append("💻 Technical/programming content")
        if any(word in text_lower for word in ['price', 'buy', 'product', 'service']):
            content_types.append("💰 Commercial/business content")
        if re.search(r'\d{4}-\d{2}-\d{2}', text):
            content_types.append("📅 Contains dates")
        
        if content_types:
            for ct in content_types:
                st.write(ct)
        else:
            st.write("📄 General text content")
    
    elif st.session_state.current_task == "search":
        st.markdown("### 🔍 Search Within Content")
        
        search_query = st.text_input("Enter keyword or phrase:", placeholder="Type to search...")
        
        if search_query:
            text_lower = st.session_state.content_text.lower()
            query_lower = search_query.lower()
            
            # Count occurrences
            count = text_lower.count(query_lower)
            
            if count > 0:
                st.success(f"✅ Found '{search_query}' {count} time(s)!")
                
                # Find and show context
                st.markdown("**Occurrences with context:**")
                
                # Split into sentences
                sentences = re.split(r'[.!?]+', st.session_state.content_text)
                
                matches = []
                for sent in sentences:
                    if query_lower in sent.lower():
                        matches.append(sent.strip())
                
                for i, match in enumerate(matches[:10], 1):
                    # Highlight the search term
                    highlighted = re.sub(
                        f'({re.escape(search_query)})',
                        r'**\1**',
                        match,
                        flags=re.IGNORECASE
                    )
                    st.markdown(f"{i}. {highlighted}...")
            else:
                st.warning(f"❌ '{search_query}' not found in the content")
                
                # Suggest similar words
                words_in_text = set(st.session_state.content_text.lower().split())
                query_words = query_lower.split()
                
                similar = [w for w in words_in_text if any(qw in w for qw in query_words) and len(w) > 3]
                if similar:
                    st.info(f"💡 Did you mean: {', '.join(similar[:5])}?")
    
    elif st.session_state.current_task == "read":
        st.markdown("### 📖 Full Content")
        
        # Show in expandable sections
        st.download_button(
            label="📥 Download Full Text",
            data=st.session_state.content_text,
            file_name=f"{st.session_state.source_name}_full_text.txt",
            mime="text/plain"
        )
        
        st.text_area("Content Preview", st.session_state.content_text, height=500)

    # CUSTOM QUESTION SECTION
    st.divider()
    st.markdown("### ❓ Ask a Specific Question")
    
    question = st.text_input(
        "Your question:", 
        placeholder="e.g., What is the main topic? What are the conclusions? Explain the methodology..."
    )
    
    if question:
        st.info("💡 **How I'm answering:** I'm searching for relevant sentences that contain keywords from your question.")
        
        # Extract keywords from question
        question_words = set(question.lower().split())
        stop_words = {'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 'is', 'are', 'was', 'were', 'be', 'been', 'being', 'have', 'has', 'had', 'do', 'does', 'did', 'will', 'would', 'could', 'should', 'may', 'might', 'must', 'shall', 'can', 'need', 'what', 'which', 'who', 'when', 'where', 'why', 'how', 'this', 'that', 'these', 'those', 'it', 'its'}
        
        keywords = [w for w in question_words if w not in stop_words and len(w) > 3]
        
        if keywords:
            st.markdown(f"**Searching for:** {', '.join(keywords)}")
            
            # Find relevant sentences
            sentences = re.split(r'[.!?]+', st.session_state.content_text)
            relevant = []
            
            for sent in sentences:
                sent_lower = sent.lower()
                # Score by keyword matches
                score = sum(1 for kw in keywords if kw in sent_lower)
                if score > 0:
                    relevant.append((sent.strip(), score))
            
            relevant.sort(key=lambda x: x[1], reverse=True)
            
            if relevant:
                st.success("📌 **Relevant Information Found:**")
                for sent, score in relevant[:5]:
                    st.markdown(f"{'✅' if score >= 2 else '•'} {sent}.")
            else:
                st.warning("No exact matches found. Try using different keywords or use the Search tab.")
        else:
            st.warning("Please enter a more specific question with meaningful keywords.")

else:
    # Welcome screen
    st.divider()
    st.markdown("""
    ### 🎯 What I Can Do:
    
    **📁 File Types Supported:**
    - PDF documents
    - Word documents (DOCX)
    - PowerPoint presentations (PPTX)
    - Text files (TXT, MD)
    - Data files (CSV, JSON)
    
    **🌐 Or Paste Any Website URL**
    
    **🔧 Features:**
    1. **📝 Summarize** - Get executive summary and key points
    2. **🔑 Extract Key Points** - Pull out bullet points and important statements
    3. **📊 Analyze** - Statistics, keywords, content type detection
    4. **🔍 Search** - Find specific terms with context
    5. **📖 Read All** - View/download full content
    6. **❓ Q&A** - Ask specific questions about the content
    
    **💡 Examples:**
    - Upload a research paper → Get summary and key findings
    - Paste a news article → Analyze main topics and sentiment
    - Upload a report → Extract action items and conclusions
    - Paste a blog → Get key takeaways
    """)

# Footer
st.divider()
st.markdown("""
<center>
Built with ❤️ by Ayushi | Like ChatGPT but for YOUR documents
</center>
""", unsafe_allow_html=True)
